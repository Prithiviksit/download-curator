"""
Office documents metadata extractors: DOCX, XLSX, PPTX.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import List, Set

from download_curator.core.models import ExtractedMetadata
from download_curator.extractors.base import BaseExtractor


class DocxExtractor(BaseExtractor):
    @property
    def supported_extensions(self) -> Set[str]:
        return {".docx", ".doc"}

    def extract(self, file_path: Path) -> ExtractedMetadata:
        metadata = ExtractedMetadata(file_type="docx")
        if file_path.suffix.lower() == ".doc":
            metadata.file_type = "doc"
            return metadata

        try:
            import docx
            doc = docx.Document(str(file_path))
            core_props = doc.core_properties
            if core_props.title:
                metadata.title = core_props.title.strip()
            if core_props.author:
                authors = [a.strip() for a in re.split(r"[,;]| and ", core_props.author) if a.strip()]
                if authors:
                    metadata.authors = authors
            if core_props.created:
                metadata.year = core_props.created.year
                metadata.date = core_props.created.strftime("%Y-%m-%d")

            # Extract first 5 paragraphs
            paragraphs_text = []
            for p in doc.paragraphs[:10]:
                text = p.text.strip()
                if text:
                    paragraphs_text.append(text)

            if paragraphs_text:
                full_text = " ".join(paragraphs_text)
                metadata.excerpt = full_text[:1000]

                # If title not in properties, use first heading/paragraph
                if not metadata.title and paragraphs_text:
                    metadata.title = paragraphs_text[0][:80]

        except Exception as e:
            metadata.raw_metadata["error"] = str(e)

        return metadata


class XlsxExtractor(BaseExtractor):
    @property
    def supported_extensions(self) -> Set[str]:
        return {".xlsx", ".xls"}

    def extract(self, file_path: Path) -> ExtractedMetadata:
        metadata = ExtractedMetadata(file_type="xlsx")
        if file_path.suffix.lower() == ".xls":
            metadata.file_type = "xls"
            return metadata

        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            sheet_names = wb.sheetnames
            metadata.raw_metadata["sheets"] = sheet_names

            if sheet_names:
                first_sheet = wb[sheet_names[0]]
                headers = []
                for row in first_sheet.iter_rows(max_row=3, values_only=True):
                    row_vals = [str(v).strip() for v in row if v is not None]
                    if row_vals and not headers:
                        headers = row_vals
                metadata.raw_metadata["headers"] = headers[:15]
                metadata.excerpt = f"Sheets: {', '.join(sheet_names[:5])}. Columns: {', '.join(headers[:10])}"

            wb.close()
        except Exception as e:
            metadata.raw_metadata["error"] = str(e)

        return metadata


class PptxExtractor(BaseExtractor):
    @property
    def supported_extensions(self) -> Set[str]:
        return {".pptx", ".ppt"}

    def extract(self, file_path: Path) -> ExtractedMetadata:
        metadata = ExtractedMetadata(file_type="pptx")
        if file_path.suffix.lower() == ".ppt":
            metadata.file_type = "ppt"
            return metadata

        try:
            import pptx
            prs = pptx.Presentation(str(file_path))
            core_props = prs.core_properties
            if core_props.title:
                metadata.title = core_props.title.strip()
            if core_props.author:
                metadata.authors = [core_props.author.strip()]

            # Extract text from first 3 slides
            slides_text = []
            for slide in list(prs.slides)[:3]:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            t = paragraph.text.strip()
                            if t:
                                slides_text.append(t)

            if slides_text:
                metadata.excerpt = " | ".join(slides_text[:10])[:1000]
                if not metadata.title and slides_text:
                    metadata.title = slides_text[0][:80]
                    metadata.topic_or_subject = slides_text[0][:60]

        except Exception as e:
            metadata.raw_metadata["error"] = str(e)

        return metadata
