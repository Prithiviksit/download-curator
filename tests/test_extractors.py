"""
Unit tests for content and metadata extractors.
"""

from __future__ import annotations

import zipfile
from pathlib import Path
from PIL import Image

from download_curator.extractors.archive import ArchiveExtractor
from download_curator.extractors.image import ImageExtractor
from download_curator.extractors.installer import InstallerExtractor
from download_curator.extractors.office import DocxExtractor, PptxExtractor, XlsxExtractor
from download_curator.extractors.registry import extract_metadata
from download_curator.extractors.text import CsvExtractor, TextExtractor


def test_text_extractor_markdown(tmp_path: Path) -> None:
    md_file = tmp_path / "notes.md"
    md_file.write_text("# Meeting Notes\n\nDiscussion about architecture and rollout.\n")

    extractor = TextExtractor()
    assert extractor.can_handle(md_file)

    meta = extractor.extract(md_file)
    assert meta.title == "Meeting Notes"
    assert "Discussion about architecture" in (meta.excerpt or "")


def test_text_extractor_code(tmp_path: Path) -> None:
    py_file = tmp_path / "analyzer.py"
    py_file.write_text('"""Data processing engine."""\n\ndef run():\n    pass\n')

    extractor = TextExtractor()
    meta = extractor.extract(py_file)
    assert meta.file_type == "python"
    assert "Data processing engine" in (meta.title or "")


def test_csv_extractor(tmp_path: Path) -> None:
    csv_file = tmp_path / "sales_2026.csv"
    csv_file.write_text("date,merchant,amount,category\n2026-08-01,Acme,49.99,Software\n")

    extractor = CsvExtractor()
    meta = extractor.extract(csv_file)
    assert meta.file_type == "csv"
    assert "date" in meta.raw_metadata.get("headers", [])
    assert meta.dataset_name == "sales_2026"


def test_image_extractor(tmp_path: Path) -> None:
    img_file = tmp_path / "test_photo.png"
    # Create simple image with PIL
    img = Image.new("RGB", (100, 100), color="blue")
    img.save(img_file)

    extractor = ImageExtractor()
    assert extractor.can_handle(img_file)
    meta = extractor.extract(img_file)
    assert meta.raw_metadata.get("width") == 100
    assert meta.raw_metadata.get("height") == 100


def test_archive_extractor(tmp_path: Path) -> None:
    zip_file = tmp_path / "project_src.zip"
    with zipfile.ZipFile(zip_file, "w") as zf:
        zf.writestr("src/main.py", "print('hello')")
        zf.writestr("src/utils.py", "pass")
        zf.writestr("README.md", "# Project")

    extractor = ArchiveExtractor()
    meta = extractor.extract(zip_file)
    assert meta.file_type == "archive"
    assert meta.raw_metadata.get("total_files") == 3


def test_installer_extractor() -> None:
    dmg_path = Path("/mock/Downloads/Discord-0.0.300-arm64.dmg")
    extractor = InstallerExtractor()
    meta = extractor.extract(dmg_path)
    assert meta.application_name == "Discord"
    assert meta.version == "0.0.300"
    assert meta.architecture == "arm64"


def test_docx_extractor(tmp_path: Path) -> None:
    import docx

    doc_file = tmp_path / "report.docx"
    doc = docx.Document()
    doc.core_properties.title = "Quarterly Financial Analysis"
    doc.core_properties.author = "Jane Doe"
    doc.add_paragraph("Introduction to quarterly performance.")
    doc.save(str(doc_file))

    extractor = DocxExtractor()
    meta = extractor.extract(doc_file)
    assert meta.title == "Quarterly Financial Analysis"
    assert "Jane Doe" in meta.authors
    assert "Introduction to quarterly performance" in (meta.excerpt or "")


def test_xlsx_extractor(tmp_path: Path) -> None:
    import openpyxl

    xlsx_file = tmp_path / "dataset.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Summary"
    ws.append(["ID", "Name", "Score"])
    ws.append([1, "Alice", 95])
    wb.save(str(xlsx_file))

    extractor = XlsxExtractor()
    meta = extractor.extract(xlsx_file)
    assert "Summary" in meta.raw_metadata.get("sheets", [])
    assert meta.raw_metadata.get("headers") == ["ID", "Name", "Score"]


def test_pptx_extractor(tmp_path: Path) -> None:
    import pptx

    pptx_file = tmp_path / "presentation.pptx"
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Deep Learning Architecture"
    prs.save(str(pptx_file))

    extractor = PptxExtractor()
    meta = extractor.extract(pptx_file)
    assert "Deep Learning Architecture" in (meta.title or "")
